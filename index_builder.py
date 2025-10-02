# index_builder.py
import os, sqlite3, pathlib
from pptx import Presentation  # pip install python-pptx

def ensure_schema(dbfile):
    con = sqlite3.connect(dbfile)
    c = con.cursor()
    c.execute("CREATE TABLE IF NOT EXISTS docs (id INTEGER PRIMARY KEY, path TEXT, title TEXT, body TEXT)")
    c.execute("CREATE VIRTUAL TABLE IF NOT EXISTS docs_fts USING fts5(body, content='docs', content_rowid='id')")
    con.commit(); con.close()

def add_doc(dbfile, path, title, body):
    con = sqlite3.connect(dbfile)
    c = con.cursor()
    c.execute("INSERT INTO docs(path,title,body) VALUES (?,?,?)", (path, title, body))
    rowid = c.lastrowid
    c.execute("INSERT INTO docs_fts(rowid, body) VALUES (?,?)", (rowid, body))
    con.commit(); con.close()

def extract_text_from_pptx(pth):
    prs = Presentation(pth)
    txt = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt.append(shape.text)
    return "\n".join(txt)

def walk_and_index(dbfile, root, kind="text"):
    ensure_schema(dbfile)
    for dirpath, _, filenames in os.walk(root):
        for f in filenames:
            p = pathlib.Path(dirpath, f)
            text, title = "", p.stem
            try:
                if kind=="ppt" and p.suffix.lower()==".pptx":
                    text = extract_text_from_pptx(p)
                elif kind in ("text","media"):
                    if p.suffix.lower() in [".txt",".vtt"]:
                        text = p.read_text(encoding="utf-8", errors="ignore")
                if text.strip():
                    add_doc(dbfile, str(p), title, text)
            except Exception as e:
                print("skip", p, e)
